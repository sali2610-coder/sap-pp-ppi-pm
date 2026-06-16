#!/usr/bin/env python3
"""Presentation Review Board — score a .pptx on 10 dimensions against a template profile.

Usage:
    python3 review.py <deck.pptx> --template sap|university|executive|financial|technical_training [--out dir]

Extracts measurable signals (action titles, charts, tables, infographic images, text density, citations,
references slide, decision/action slide, native animations, font sizes), scores each dimension 1–5 from
those signals, applies the template's weights for a weighted overall, and writes:
    scorecard.json  + review.md  (scorecard · strengths · weaknesses · improvement plan · slide-by-slide)

Deterministic. Heuristic proxies make the scoring transparent and reproducible; a human reviewer can
override any dimension. Pure stdlib + python-pptx.
"""
import sys, re, json, zipfile, argparse, pathlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

HERE = pathlib.Path(__file__).resolve().parent
DECISION = ("action", "decision", "conclusion", "recommend", "next step")
SENT = re.compile(r"[a-z].*\b(is|are|holds|shows|predicts|grows|moves|ready|coincides|gate|carries|follows|runs|means|earn|fall|rises|up|down|—|:)\b", re.I)

def slide_text(s):
    return " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
def first_title(s):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""
def min_body_pt(prs):
    sizes = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size: sizes.append(int(r.font.size.pt))
    return min(sizes) if sizes else None

def clamp(x): return max(1, min(5, x))

def extract(path):
    prs = Presentation(path)
    slides = list(prs.slides)
    per = []
    for i, s in enumerate(slides, 1):
        txt = slide_text(s); title = first_title(s)
        ch = sum(1 for sh in s.shapes if sh.has_chart)
        tb = sum(1 for sh in s.shapes if sh.has_table)
        im = sum(1 for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
        words = len(txt.split())
        cite = bool(re.search(r"source|references|\(20\d\d\)|et al", txt, re.I))
        action = bool(SENT.search(title)) and len(title.split()) >= 4
        decision = any(k in txt.lower() for k in DECISION)
        per.append({"n": i, "title": title, "words": words, "charts": ch, "tables": tb,
                    "images": im, "citation": cite, "action_title": action, "decision": decision})
    # native animation signals from the zip
    tim = tr = 0
    with zipfile.ZipFile(path) as z:
        for nm in z.namelist():
            if re.match(r"ppt/slides/slide\d+\.xml$", nm):
                b = z.read(nm); tim += b.count(b"<p:timing>"); tr += b.count(b"<p:transition")
    return prs, slides, per, {"timing": tim, "transition": tr, "min_body_pt": min_body_pt(prs)}

def score(per, agg, tpl):
    n = len(per)
    charts = sum(p["charts"] for p in per); tables = sum(p["tables"] for p in per)
    infog = sum(p["images"] for p in per)
    avg_words = sum(p["words"] for p in per) / max(1, n)
    cited = sum(p["citation"] for p in per); refs = any("reference" in p["title"].lower() for p in per)
    action_ratio = sum(p["action_title"] for p in per[1:]) / max(1, n - 1)  # exclude title slide
    decision = any(p["decision"] for p in per)
    decision_early = any(p["decision"] for p in per[:max(2, n // 2)])
    has_motiv = n >= 4
    minpt = agg["min_body_pt"]; tgt = tpl.get("min_body_pt", 18); wtgt = tpl.get("body_word_target", 45)
    D = {}
    # 1 executive
    D["executive_quality"] = clamp(round(2 + (1 if decision_early else 0) + (1.5 if avg_words <= wtgt else 0) + (0.5 if infog or charts else 0)))
    # 2 academic
    D["academic_quality"] = clamp(round(2 + (1.5 if refs else 0) + (1.5 if cited >= 1 else 0)))
    # 3 storytelling
    D["storytelling"] = clamp(round(2 + 2 * action_ratio + (1 if (decision and has_motiv) else 0)))
    # 4 visual design
    textonly = sum(1 for p in per if p["charts"] + p["tables"] + p["images"] == 0 and p["words"] > 30)
    D["visual_design"] = clamp(round(5 - min(3, textonly)))
    # 5 data viz
    D["data_visualization"] = clamp(round(2 + (2 if charts >= 1 else 0) + (1 if (charts + tables) >= 2 else 0)))
    # 6 infographics
    D["infographics"] = clamp(round(2 + (2 if infog >= 1 else 0) + (1 if infog >= 2 else 0)))
    # 7 animation
    D["animation_usage"] = clamp(round(2 + (2 if agg["timing"] >= n else (1 if agg["timing"] else 0)) + (1 if agg["transition"] >= n else 0)))
    # 8 accessibility
    D["accessibility"] = clamp(round(2 + (2 if (minpt is None or minpt >= tgt) else 0) + (1 if (minpt is None or minpt >= 16) else 0)))
    # 9 audience fit (density vs target + has structure)
    D["audience_fit"] = clamp(round(3 + (1 if avg_words <= wtgt else -1) + (1 if (charts and infog) else 0)))
    # 10 decision readiness
    D["decision_readiness"] = clamp(round(2 + (2 if decision else 0) + (1 if per[-1]["decision"] else 0)))
    facts = {"slides": n, "charts": charts, "tables": tables, "infographic_images": infog,
             "avg_words": round(avg_words, 1), "citations": cited, "references_slide": refs,
             "action_title_ratio": round(action_ratio, 2), "decision_slide": decision,
             "timing_slides": agg["timing"], "transition_slides": agg["transition"], "min_body_pt": minpt}
    return D, facts

def weighted(D, tpl):
    w = tpl["weights"]
    return round(sum(D[k] * w[k] for k in D), 2)

def narrative(D, facts, tpl, per):
    strengths, weaknesses, plan = [], [], []
    for k, v in sorted(D.items(), key=lambda x: -x[1]):
        label = k.replace("_", " ")
        if v >= 4: strengths.append(f"{label} ({v}/5)")
        elif v <= 2: weaknesses.append(f"{label} ({v}/5)")
    # targeted improvements
    if D["infographics"] <= 3: plan.append("Add a process/relationship infographic (use infographics/ library) on a structural slide.")
    if D["data_visualization"] <= 3: plan.append("Add a native chart on each data slide; pick type via chart-engine.")
    if D["animation_usage"] <= 3: plan.append("Run upgrades/pptx_animate.py to add native entrance timing + transitions.")
    if D["storytelling"] <= 3: plan.append("Rewrite topic-label titles as full-sentence action titles; verify the ghost-deck.")
    if D["academic_quality"] <= 3 and tpl["profile"] in ("University",): plan.append("Add in-slide citations and a References slide.")
    if D["accessibility"] <= 3: plan.append(f"Raise body text to >= {tpl.get('min_body_pt',18)}pt; add image alt text.")
    if D["decision_readiness"] <= 3: plan.append("End on an explicit decision/action slide with a clear recommendation.")
    if not plan: plan.append("No structural gaps; polish wording and annotate the key finding on each chart.")
    return strengths, weaknesses, plan

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck"); ap.add_argument("--template", required=True)
    ap.add_argument("--out", default=None)
    a = ap.parse_args()
    tpl = json.load(open(HERE / "templates" / f"{a.template}.json"))
    prs, slides, per, agg = extract(a.deck)
    D, facts = score(per, agg, tpl)
    overall = weighted(D, tpl)
    strengths, weaknesses, plan = narrative(D, facts, tpl, per)
    out = pathlib.Path(a.out) if a.out else pathlib.Path(a.deck).parent
    out.mkdir(parents=True, exist_ok=True)
    scorecard = {"deck": str(a.deck), "template": tpl["profile"], "overall_weighted": overall,
                 "dimensions": D, "facts": facts}
    (out / "scorecard.json").write_text(json.dumps(scorecard, indent=2))
    # markdown report
    L = [f"# Review — {pathlib.Path(a.deck).parent.name}", f"Template: **{tpl['profile']}** · Overall: **{overall}/5**\n",
         "## Scorecard", "| Dimension | Score | Weight |", "|---|:--:|:--:|"]
    for k in D: L.append(f"| {k.replace('_',' ')} | {D[k]}/5 | {tpl['weights'][k]:.2f} |")
    str_lines = [f"- {x}" for x in strengths] or ["- (none ≥4)"]
    weak_lines = [f"- {x}" for x in weaknesses] or ["- none ≤2"]
    plan_lines = [f"{i+1}. {x}" for i, x in enumerate(plan)]
    L += ["", "## Strengths", *str_lines,
          "", "## Weaknesses", *weak_lines,
          "", "## Improvement plan", *plan_lines,
          "", "## Slide-by-slide"]
    for p in per:
        tags = []
        if p["action_title"]: tags.append("action-title")
        if p["charts"]: tags.append(f"{p['charts']} chart")
        if p["tables"]: tags.append(f"{p['tables']} table")
        if p["images"]: tags.append(f"{p['images']} infographic")
        if p["citation"]: tags.append("cited")
        if p["decision"]: tags.append("decision")
        flag = "⚠ text-heavy" if (p["words"] > tpl.get("body_word_target", 45) and not (p["charts"] or p["images"] or p["tables"])) else ""
        L.append(f"- **S{p['n']}** {p['title'][:70]} — {', '.join(tags) or 'text'} ({p['words']}w) {flag}")
    (out / "review.md").write_text("\n".join(L) + "\n")
    print(json.dumps({"template": tpl["profile"], "overall": overall, "dimensions": D}, indent=2))

if __name__ == "__main__":
    main()
