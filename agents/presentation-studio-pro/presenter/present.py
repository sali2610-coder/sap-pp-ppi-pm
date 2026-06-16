#!/usr/bin/env python3
"""AI Presenter — turn a .pptx into a delivery kit for a given audience.

Usage:
    python3 present.py <deck.pptx> --profile academic|sap|executive|training|financial [--minutes N] [--out dir]

Generates (presenter_kit.md + timing.json):
  1. Speaker notes (per slide)      4. Audience objections (+ rebuttals)
  2. Presentation script (narration) 5. Timing plan (per slide + total vs target)
  3. Q&A preparation (questions+ans) 6. Executive summary

Reads slide titles + bullets, computes timing from word counts and the profile's words-per-minute, and
scaffolds notes/script/Q&A/objections from the actual deck content + the audience profile. Deterministic.
The agent refines the scaffold; the timing + structure are exact.
"""
import sys, re, json, argparse, pathlib
from pptx import Presentation

HERE = pathlib.Path(__file__).resolve().parent

def slides_data(path):
    prs = Presentation(path)
    out = []
    for i, s in enumerate(prs.slides, 1):
        title, body = "", []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                lines = [l.strip() for l in sh.text_frame.text.split("\n") if l.strip()]
                if not title:
                    title = lines[0]; body += lines[1:]
                else:
                    body += lines
        words = sum(len(b.split()) for b in body) + len(title.split())
        kind = "title" if i == 1 else ("decision" if re.search(r"action|decision|conclusion|recommend", title, re.I) else
               ("result" if re.search(r"\d|%|×|grow|fall|predict|coincid|gap", title) else "content"))
        out.append({"n": i, "title": title, "body": body, "words": words, "kind": kind})
    return out

def speaker_notes(sl, prof):
    cue = prof["notes_style"]
    notes = []
    for s in sl:
        pts = "; ".join(s["body"][:4]) if s["body"] else "(visual / chart carries the point)"
        notes.append(f"**S{s['n']} — {s['title']}**\n  Cue ({prof['tone']}): {cue}.\n  Say: {pts}\n  Bridge: \"which leads to…\" → next slide.")
    return "\n".join(notes)

def script(sl):
    out = []
    for s in sl:
        if s["kind"] == "title":
            out.append(f"[Open] \"{s['title']}.\" Set context in one breath, then preview the through-line.")
        elif s["kind"] == "decision":
            out.append(f"[Close] \"{s['title']}.\" State the ask plainly and stop talking.")
        else:
            lead = s["title"]
            ev = (" " + " ".join(s["body"][:2])) if s["body"] else ""
            out.append(f"\"{lead}.\"{(' ' + ev) if ev else ''} Land the takeaway, then move on.")
    return "\n\n".join(f"{i+1}. {p}" for i, p in enumerate(out))

def qa(sl, prof):
    qs = []
    # profile-driven
    for f in prof["qa_focus"]:
        qs.append((f"On {f}: …?", f"Anchor to the evidence on the relevant slide; concede the limit, then give the strongest defensible point."))
    # slide-derived (results/decision)
    for s in sl:
        if s["kind"] == "result":
            qs.append((f"How robust is \"{s['title']}\"?",
                       "Give the number's source + sensitivity; note what would change it."))
        if s["kind"] == "decision":
            qs.append((f"What do you need from us on \"{s['title']}\"?",
                       "State the specific decision, owner, and date."))
    seen = set(); kit = []
    for q, a in qs:
        if q in seen: continue
        seen.add(q); kit.append(f"- **Q:** {q}\n  **A:** {a}")
    return "\n".join(kit[:10])

def objections(prof):
    return "\n".join(f"- **\"{o}\"** → Rebuttal: acknowledge, then counter with the deck's evidence; if valid, name the mitigation."
                     for o in prof["objections"])

def timing(sl, prof, target_min):
    wpm = prof["wpm"]; rows = []; total = 0.0
    for s in sl:
        spoken = s["words"] / wpm
        # add base time for the slide (setup/transition/visual dwell)
        base = 0.6 if s["kind"] in ("result", "decision") else 0.4
        mins = round(spoken + base, 2); total += mins
        rows.append({"slide": s["n"], "title": s["title"][:50], "minutes": mins})
    total = round(total, 2)
    verdict = "on target"
    if target_min:
        if total > target_min * 1.05: verdict = f"OVER by {round(total-target_min,1)} min — cut content/words"
        elif total < target_min * 0.85: verdict = f"UNDER by {round(target_min-total,1)} min — add depth or Q&A"
    return {"wpm": wpm, "target_min": target_min, "estimated_min": total, "verdict": verdict, "per_slide": rows}

def exec_summary(sl, prof):
    title = sl[0]["title"] if sl else ""
    results = [s["title"] for s in sl if s["kind"] == "result"][:2]
    decision = next((s["title"] for s in sl if s["kind"] == "decision"), "")
    parts = [f"{title}."]
    if results: parts.append("Key findings: " + "; ".join(results) + ".")
    if decision: parts.append("Recommendation: " + decision + ".")
    text = " ".join(parts)
    # trim toward profile.summary_words (soft)
    return text

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck"); ap.add_argument("--profile", required=True)
    ap.add_argument("--minutes", type=float, default=None); ap.add_argument("--out", default=None)
    a = ap.parse_args()
    prof = json.load(open(HERE / "profiles.json"))[a.profile]
    sl = slides_data(a.deck)
    tm = timing(sl, prof, a.minutes)
    out = pathlib.Path(a.out) if a.out else pathlib.Path(a.deck).parent
    out.mkdir(parents=True, exist_ok=True)
    (out / "timing.json").write_text(json.dumps(tm, indent=2))
    md = []
    md.append(f"# Presenter Kit — {pathlib.Path(a.deck).parent.name}")
    md.append(f"Audience profile: **{a.profile}** · tone: {prof['tone']} · {prof['wpm']} wpm\n")
    md.append("## 6. Executive summary\n" + exec_summary(sl, prof) + "\n")
    md.append(f"## 5. Timing plan\nEstimated **{tm['estimated_min']} min**" +
              (f" vs target {a.minutes} min → **{tm['verdict']}**" if a.minutes else "") + "\n")
    md.append("| Slide | Title | Min |\n|--|--|--|")
    for r in tm["per_slide"]: md.append(f"| {r['slide']} | {r['title']} | {r['minutes']} |")
    md.append(f"| **Total** |  | **{tm['estimated_min']}** |\n")
    md.append("## 1. Speaker notes\n" + speaker_notes(sl, prof) + "\n")
    md.append("## 2. Presentation script\n" + script(sl) + "\n")
    md.append("## 3. Q&A preparation\n" + qa(sl, prof) + "\n")
    md.append("## 4. Audience objections\n" + objections(prof) + "\n")
    (out / "presenter_kit.md").write_text("\n".join(md) + "\n")
    print(json.dumps({"profile": a.profile, "slides": len(sl),
                      "estimated_min": tm["estimated_min"], "target": a.minutes,
                      "verdict": tm["verdict"]}, indent=2))

if __name__ == "__main__":
    main()
