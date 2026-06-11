import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

D=json.load(open('exports/sap-architecture-data.json'))
PAL=D['palette']
RED=RGBColor(0xd6,0x20,0x27); GOLD=RGBColor(0xf2,0xc1,0x4e); INK=RGBColor(0x07,0x0b,0x11)
PANEL=RGBColor(0x0e,0x16,0x20); TXT=RGBColor(0xea,0xf0,0xf7); MUT=RGBColor(0x90,0xa0,0xb3)
LINE=RGBColor(0x1d,0x2a,0x38); WHITE=RGBColor(0xff,0xff,0xff)
def hexc(h): h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]

def bg(s,color=INK):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=color
def rect(s,x,y,w,h,fill,line=None,lw=1.0,rad=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,x,y,w,h)
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    return shp
def rtl(p):
    pPr=p._p.get_or_add_pPr(); pPr.set('rtl','1')
def txt(s,x,y,w,h,text,size=18,color=TXT,bold=False,align=PP_ALIGN.RIGHT,rtlf=True,font='Segoe UI',anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    lines=text.split('\n')
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if rtlf: rtl(p)
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=font
    return tb
def chip(s,x,y,w,v,label):
    rect(s,x,y,w,Inches(0.95),PANEL,LINE,1,True)
    txt(s,x,y+Inches(0.08),w,Inches(0.5),str(v),28,WHITE,True,PP_ALIGN.CENTER,False,'Consolas')
    txt(s,x,y+Inches(0.6),w,Inches(0.3),label,10,MUT,False,PP_ALIGN.CENTER,True)

# accent bar
def accent(s):
    rect(s,0,0,SW,Pt(5),RED)

# ---------- 1 TITLE ----------
s=prs.slides.add_slide(blank); bg(s)
rect(s,Inches(0.9),Inches(2.5),Inches(0.9),Inches(0.9),RED,None,0,True)
txt(s,Inches(0.95),Inches(2.62),Inches(0.8),Inches(0.7),'NEO',20,WHITE,True,PP_ALIGN.CENTER,False)
txt(s,Inches(2.0),Inches(2.4),Inches(10.5),Inches(1.0),'מפת ארכיטקטורת SAP — חדר בקרה',40,WHITE,True)
txt(s,Inches(2.0),Inches(3.5),Inches(10.5),Inches(0.6),'נוף SAP מלא · ECC6 → S/4HANA · CBC Israel',20,GOLD,False)
txt(s,Inches(2.0),Inches(4.2),Inches(10.5),Inches(0.5),'PM · PP-PI · MM · SD · FI · CO · QM · CS · Integration',14,MUT,False,PP_ALIGN.RIGHT,False,'Consolas')
txt(s,Inches(0.9),Inches(6.7),Inches(11.5),Inches(0.4),'נבנה ע״י Sali Halif — Web Coding · 2026',12,MUT,False,PP_ALIGN.RIGHT)
accent(s)
c=D['meta']['counts']
xs=Inches(8.7)
for i,(v,l) in enumerate([(c['tables'],'טבלאות'),(c['modules'],'מודולים'),(c['shared'],'צמתים'),(c['edges'],'קשרים')]):
    chip(s,xs+Inches(i*1.15),Inches(0.4),Inches(1.05),v,l)

def header(s,he,en):
    accent(s); txt(s,Inches(0.6),Inches(0.25),Inches(9),Inches(0.7),he,26,WHITE,True)
    txt(s,Inches(0.6),Inches(0.95),Inches(9),Inches(0.4),en,13,GOLD,False,PP_ALIGN.RIGHT,False,'Consolas')

def img_full(s,path,x,y,w):
    try: s.shapes.add_picture(path,x,y,width=w)
    except Exception as e: print('img',path,e)

# ---------- 2 LANDSCAPE ----------
s=prs.slides.add_slide(blank); bg(s); header(s,'נוף ה-SAP המלא — חמש שכבות','Five-Layer Landscape')
img_full(s,'exports/slide-landscape.png',Inches(0.6),Inches(1.5),Inches(8.5))
for i,t in enumerate(['שכבה 1 · תהליכים עסקיים (P2P · O2C · M2O · S2C · Plan2Produce)',
   'שכבה 2 · 9 מודולי SAP','שכבה 3 · 13 מסמכים עסקיים','שכבה 4 · 134 טבלאות מסד',
   'שכבה 5 · 14 אובייקטי ליבה משותפים — במרכז']):
    txt(s,Inches(9.3),Inches(1.7+i*0.62),Inches(3.6),Inches(0.6),'• '+t,12.5,TXT,i==4,PP_ALIGN.RIGHT)

# ---------- 3 PROCESSES ----------
s=prs.slides.add_slide(blank); bg(s); header(s,'תהליכים עסקיים מקצה-לקצה','End-to-End Business Processes')
img_full(s,'exports/slide-business.png',Inches(0.6),Inches(1.5),Inches(8.4))
y=Inches(1.7)
for p in D['processes']:
    rect(s,Inches(9.2),y,Inches(0.18),Inches(0.78),hexc(p['color']))
    txt(s,Inches(9.5),y,Inches(3.4),Inches(0.4),p['name'],14,WHITE,True)
    txt(s,Inches(9.5),y+Inches(0.36),Inches(3.4),Inches(0.4),p['he']+' · '+' · '.join(p['mods']),10.5,MUT,False,PP_ALIGN.RIGHT,True,'Segoe UI')
    y+=Inches(1.02)

# ---------- 4 SHARED HUBS (table) ----------
s=prs.slides.add_slide(blank); bg(s); header(s,'אובייקטים משותפים — צמתים מרכזיים','Shared Core Objects · Hubs')
hubs=sorted([h for h in D['shared'] if h['degree'] is not None],key=lambda x:-x['degree'])+[h for h in D['shared'] if h['degree'] is None]
rows=len(hubs)+1
tb=s.shapes.add_table(rows,5,Inches(0.6),Inches(1.5),Inches(8.3),Inches(5.6)).table
for j,h in enumerate(['טבלה','קשרים','אב','צפיפות','מודולים']):
    cell=tb.cell(0,j); cell.text=h; cell.fill.solid(); cell.fill.fore_color.rgb=RED
    pp=cell.text_frame.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    pp.runs[0].font.size=Pt(11); pp.runs[0].font.bold=True; pp.runs[0].font.color.rgb=WHITE
for i,h in enumerate(hubs,1):
    vals=[h['name'],str(h['degree'] if h['degree'] is not None else '—'),str(h['inDegree'] if h['inDegree'] is not None else '—'),
          str(h['density'] if h['density'] is not None else 'land'),str(h['span'])+' · '+','.join(h['mods'][:4])]
    for j,v in enumerate(vals):
        cell=tb.cell(i,j); cell.text=v; cell.fill.solid(); cell.fill.fore_color.rgb=PANEL if i%2 else RGBColor(0x0b,0x12,0x1a)
        r=cell.text_frame.paragraphs[0].runs[0]; r.font.size=Pt(10); r.font.name='Consolas'
        r.font.color.rgb=GOLD if j==0 else TXT; r.font.bold=(j==0)
        cell.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER if j<4 else PP_ALIGN.RIGHT
txt(s,Inches(9.2),Inches(1.6),Inches(3.7),Inches(5),'MARA ו-AUFK הם הצמתים המרכזיים של הנוף.\n\nאובייקטים משותפים אינם שייכים למודול בודד — הם עמוד השדרה שמחבר רכש, ייצור, תחזוקה, מכירות וכספים.\n\nגודל = מספר קשרים · צבע = צפיפות חיבור.',13,TXT,False,PP_ALIGN.RIGHT)

# ---------- 5 DATABASE ----------
s=prs.slides.add_slide(blank); bg(s); header(s,'רשת הטבלאות — תצוגת מסד נתונים','Database · Force Network')
img_full(s,'exports/slide-database.png',Inches(0.6),Inches(1.5),Inches(9.6))
txt(s,Inches(10.4),Inches(1.7),Inches(2.6),Inches(5),'134 טבלאות\n140 קשרים\n34 חוצי-מודול\n\nקשרים חוצי-מודול מודגשים בזהב — שם מתרכז הסיכון בהמרה ל-S/4HANA.',13,TXT,False,PP_ALIGN.RIGHT)

# ---------- 6 INTEGRATION ----------
s=prs.slides.add_slide(blank); bg(s); header(s,'שכבת אינטגרציה — Zetes / Daymax','Integration · IDOC / ALE / PI-PO')
img_full(s,'exports/slide-integration.png',Inches(0.6),Inches(1.5),Inches(8.0))
y=Inches(1.7)
for f in D['integrationFlows']:
    txt(s,Inches(8.8),y,Inches(4.1),Inches(0.34),f['idoc'],13,GOLD,True,PP_ALIGN.RIGHT,False,'Consolas')
    txt(s,Inches(8.8),y+Inches(0.32),Inches(4.1),Inches(0.4),f['he']+'  ('+f['from']+'→'+f['to']+')',10.5,MUT,False,PP_ALIGN.RIGHT)
    y+=Inches(0.86)

# ---------- 7 S/4 ----------
s=prs.slides.add_slide(blank); bg(s); header(s,'מעבר ECC6 → S/4HANA','S/4HANA Simplification')
img_full(s,'exports/slide-s4.png',Inches(0.6),Inches(1.5),Inches(7.2))
y=Inches(1.7)
for d in D['s4deltas']:
    txt(s,Inches(8.1),y,Inches(4.8),Inches(0.34),d['ecc']+'  →  '+d['s4'],13,WHITE,True,PP_ALIGN.RIGHT,False,'Consolas')
    txt(s,Inches(8.1),y+Inches(0.32),Inches(4.8),Inches(0.4),d['he'],10.5,MUT,False,PP_ALIGN.RIGHT)
    y+=Inches(0.92)

# ---------- 8 CLOSING ----------
s=prs.slides.add_slide(blank); bg(s); accent(s)
txt(s,Inches(1),Inches(2.6),Inches(11.3),Inches(1),'נוף SAP אחד · חוויית ידע אינטראקטיבית',32,WHITE,True,PP_ALIGN.CENTER)
txt(s,Inches(1),Inches(3.8),Inches(11.3),Inches(0.6),'7 תצוגות · 5 שכבות · 134 טבלאות · 14 צמתים · אופליין 100%',16,GOLD,False,PP_ALIGN.CENTER,False)
txt(s,Inches(1),Inches(6.7),Inches(11.3),Inches(0.4),'Sali Halif — Web Coding · NEO Cockpit · CBC Israel · 2026',12,MUT,False,PP_ALIGN.CENTER)

prs.save('exports/sap-architecture-presentation.pptx')
print('PPTX saved:',len(prs.slides.__iter__.__self__._sldIdLst),'slides')
